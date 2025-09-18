from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# Initialize presentation globally so all workers can add to the same file
prs = Presentation()

def add_slide(title: str, bullets: list[str]):
    """
    Simple tool function that adds a new slide.
    Workers will call this directly.
    - title: slide title
    - bullets: list of bullet point strings
    """
    layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(layout)

    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Content
    content_placeholder = slide.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT

def save_ppt(filepath=r"C:\\Users\\Rushil\\Desktop\\Projects\\techblocks-agenticAI-mp\\output.pptx"): # ? IDK BHAI CHAL JAA NA
    prs.save(filepath)
    print(f"[✅] PPT saved at {filepath}")